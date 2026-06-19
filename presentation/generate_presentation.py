#!/usr/bin/env python3
"""
Generator prezentacji "Inteligentne Kasyno".

Jeden model treści -> dwa artefakty:
  1. presentation.html  — deck w przeglądarce (startup style, light theme,
     fonty Inter + JetBrains Mono, akcent indygo, nawigacja strzałkami / F / P).
  2. presentation.pptx  — ta sama treść złożona natywnie przez python-pptx
     (w pełni edytowalny, wierny układowi z HTML-a).

Konwersja HTML -> PPTX:
  * Jeśli w środowisku jest playwright + Chromium, każdy slajd jest zrzucany
    do PNG i wklejany na pełny slajd (1:1 z przeglądarką).
  * W przeciwnym razie PPTX jest renderowany wektorowo z tego samego modelu
    (działa offline, bez przeglądarki) — to ścieżka domyślna tutaj.

Uruchomienie:
    python3 generate_presentation.py

Wymagania: python-pptx (pip install python-pptx). Pillow opcjonalnie.
"""

from __future__ import annotations

import os
from pathlib import Path

# ───────────────────────────── paleta ──────────────────────────────────────
INK = "0B0F1E"        # nagłówki
SLATE = "1F293A"      # treść
MUTED = "64748B"      # meta / podpisy
ACCENT = "4F46E5"     # indygo
ACCENT_SOFT = "EEF0FE"
LINE = "E6E8F0"
BG = "FFFFFF"
BG_SOFT = "F7F8FB"
GREEN = "16A34A"

BRAND = "INTELIGENTNE KASYNO"
META_TOP = "Projekt zespołowy · Inżynieria oprogramowania<br>Czerwiec 2026"

# ───────────────────────────── treść decku ─────────────────────────────────
# Każdy slajd to dict z polem "kind" i danymi. To samo zasila HTML i PPTX.
SLIDES = [
    {
        "kind": "title",
        "eyebrow": "PITCH  /  v1",
        "kicker": "Reinforcement learning · gry hazardowe online",
        "title": "Kasyno, w którym grasz<br>przeciwko agentom RL.",
        "lead": "Pełnowymiarowa platforma gier — blackjack, poker i ruletka — "
                "gdzie przeciwnikami są boty uczone metodami uczenia ze "
                "wzmocnieniem. Gra w czasie rzeczywistym, wirtualny portfel i "
                "płatności.",
        "foot": "Inteligentne Kasyno · FastAPI + React · Q-Learning & PPO · wdrożone na Render",
    },
    {
        "kind": "cards",
        "eyebrow": "01 / Problem",
        "foot": "Problem",
        "cols": 3,
        "heading": "Gry o niepełnej informacji to twardy problem.",
        "sub": "Blackjack i poker to nie czysta logika — to decyzje pod "
               "niepewnością, z ukrytą informacją i losowością.",
        "items": [
            ("01", "Ukryta informacja",
             "Karty przeciwnika i krupiera są zakryte. Decyzję podejmujesz na "
             "podstawie prawdopodobieństw, nie pewności."),
            ("02", "Losowość i wariancja",
             "Dobra decyzja może przegrać rundę. Jakość gry widać dopiero po "
             "wielu rozdaniach."),
            ("03", "Retencja gracza",
             "Kasyno musi być nie tylko sprawne, ale i wciągające — portfel, "
             "bonusy, płatności."),
        ],
        "callout": "Dlatego boty nie są skryptem „if-else”. To polityki uczone "
                   "na symulowanych rozdaniach — z gradacją trudności.",
    },
    {
        "kind": "cards",
        "eyebrow": "02 / Produkt",
        "foot": "Co zbudowaliśmy",
        "cols": 3,
        "heading": "Co zbudowaliśmy.",
        "sub": "Webowa platforma kasyna z grą w czasie rzeczywistym przeciw agentom AI.",
        "items": [
            ("01", "Trzy gry",
             "Blackjack, Texas Hold'em i ruletka europejska — stoły solo i wieloosobowe."),
            ("02", "Przeciwnicy AI",
             "Boty na trzech poziomach trudności, od losowych po uczone głęboko (PPO)."),
            ("03", "Pełny portfel",
             "Rejestracja, wirtualne żetony, historia transakcji i płatności Stripe."),
        ],
        "callout": "Frontend React + nginx, backend FastAPI, stan na żywo w "
                   "Redisie — wszystko w Dockerze, wdrożone w chmurze.",
    },
    {
        "kind": "pipeline",
        "eyebrow": "03 / Architektura",
        "foot": "Architektura",
        "heading": "Jedna pętla: akcja → stan → ruch bota → wynik.",
        "sub": "Logika gier (engine) jest odseparowana od warstwy webowej — Clean Architecture.",
        "steps": [
            ("01", "Klient", "React + TS · HTTP + WebSocket"),
            ("02", "FastAPI", "REST + WS · autoryzacja JWT"),
            ("03", "Serwisy", "Silnik gry · obsługa rundy"),
            ("04", "Redis / PG", "Gorący stan O(1) · dane ACID"),
            ("05", "Modele RL", "Inferencja ruchu bota"),
        ],
        "callout": "Redis trzyma „gorący” stan aktywnych stołów, by nie obciążać "
                   "PostgreSQL przy każdym ruchu; Postgres trzyma konta, portfele i historię.",
    },
    {
        "kind": "cards",
        "eyebrow": "04 / Technologia",
        "foot": "Stos technologiczny",
        "cols": 4,
        "heading": "Stos technologiczny.",
        "sub": None,
        "items": [
            ("BE", "Backend",
             "FastAPI (async), Python 3.11, Pydantic v2, SQLAlchemy + asyncpg, Alembic."),
            ("DB", "Dane",
             "PostgreSQL (ACID) na konta i portfele, Redis na stan stołów."),
            ("FE", "Frontend",
             "React + TypeScript + Vite, serwowane przez nginx (proxy /api i /ws)."),
            ("AI", "AI / RL",
             "Q-Learning (tabela) i PPO/DQN (Stable-Baselines3 + PyTorch)."),
        ],
        "callout": "WebSockets dają komunikację w czasie rzeczywistym; "
                   "Stripe Checkout obsługuje płatności.",
    },
    {
        "kind": "panels",
        "eyebrow": "05 / Gry",
        "foot": "Gry",
        "cols": 3,
        "heading": "Trzy gry, jeden silnik.",
        "sub": None,
        "panels": [
            ("Blackjack", None,
             ["Stoły solo i wieloosobowe z botami",
              "Pełne zasady: krupier, podział, dobór",
              "Boty od losowych po Q-Learning / PPO"]),
            ("Texas Hold'em", None,
             ["Stół 6-osobowy: blindy, fazy, showdown",
              "Stan rundy w czasie rzeczywistym",
              "Boty od luźnego callera po politykę PPO"]),
            ("Ruletka europejska", None,
             ["37 numerów, zasady europejskie",
              "Czysta gra losowa — zakłady i koło",
              "Poziom trudności jej nie zmienia"]),
        ],
        "callout": "Czysta logika gier żyje w /engine, niezależna od API i bazy "
                   "— łatwa do testów i trenowania.",
    },
    {
        "kind": "cards",
        "eyebrow": "06 / Przeciwnik AI",
        "foot": "Poziomy trudności",
        "cols": 3,
        "heading": "Trzy poziomy trudności przeciwnika.",
        "sub": "Ta sama gra, inna polityka bota — wybierana w locie z rejestru.",
        "items": [
            ("EASY", "Losowy",
             "Losowe legalne ruchy. Dla początkujących i do testów."),
            ("MEDIUM", "Heurystyka",
             "Blackjack: basic strategy. Poker: gra oparta na sile ręki."),
            ("HARD", "Deep RL",
             "Polityka PPO/DQN/Q-Learning, uczona na symulacjach. Z miękkim "
             "fallbackiem, gdy brak wag."),
        ],
        "callout": "Polityki są cache'owane per (gra, trudność) — model ładuje "
                   "się raz i jest współdzielony przez wszystkie stoły.",
    },
    {
        "kind": "cards",
        "eyebrow": "07 / Trening",
        "foot": "Jak powstają agenci",
        "cols": 3,
        "heading": "Jak powstają agenci.",
        "sub": "Osobne środowisko treningowe (/rl_training) z eksportem wag do inferencji.",
        "items": [
            ("01", "Środowiska",
             "Symulatory blackjacka i pokera w stylu Gymnasium — pełne reguły, "
             "nagroda za wynik rundy."),
            ("02", "Trening",
             "Q-Learning (tabela) dla prostych stanów, PPO/DQN (deep RL) dla pokera."),
            ("03", "Wdrożenie",
             "Wagi z /saved_models ładuje rejestr; brak torch → graceful fallback."),
        ],
        "callout": "Aplikacja nigdy nie zależy twardo od PyTorcha — gdy modelu "
                   "brak, bot schodzi do heurystyki i gra dalej.",
    },
    {
        "kind": "panels",
        "eyebrow": "08 / Real-time",
        "foot": "Czas rzeczywisty",
        "cols": 2,
        "heading": "Rozgrywka w czasie rzeczywistym.",
        "sub": None,
        "panels": [
            ("WebSockets", None,
             ["Stan stołu pushowany do graczy na żywo",
              "Autoryzacja przez jednorazowy ticket w Redisie",
              "Osobne kanały: poker oraz blackjack / ruletka"]),
            ("Gorący stan w Redisie", None,
             ["Rozdanie i stan rundy w pamięci (O(1))",
              "Bot odczytuje stan z Redisa na każdym kroku",
              "Brak obciążenia PostgreSQL przy każdym ruchu"]),
        ],
        "callout": "Setki ruchów na minutę nie dotykają bazy relacyjnej — "
                   "ta zapisuje tylko wynik rundy i portfel.",
    },
    {
        "kind": "panels",
        "eyebrow": "09 / Portfel",
        "foot": "Płatności i retencja",
        "cols": 2,
        "heading": "Portfel, płatności i retencja.",
        "sub": None,
        "panels": [
            ("Płatności — Stripe", None,
             ["Stripe Checkout na doładowanie żetonów",
              "Realizacja przez webhook checkout.session.completed",
              "Tryb symulacji bez kluczy — działa lokalnie"]),
            ("Retencja", None,
             ["Bonus powitalny przy rejestracji",
              "Bad-beat bonus za pechowe przegrane",
              "Historia transakcji i poziom bonusów w portfelu"]),
        ],
        "callout": "Portfel jest transakcyjny (ACID) — każdy zakład, wygrana i "
                   "bonus to wpis w historii.",
    },
    {
        "kind": "cards",
        "eyebrow": "10 / Wdrożenie",
        "foot": "Deployment",
        "cols": 3,
        "heading": "Od kodu do działającej chmury.",
        "sub": None,
        "items": [
            ("01", "Docker",
             "Cały stos w Docker Compose: backend, frontend, Postgres, Redis."),
            ("02", "Render",
             "Blueprint render.yaml stawia 4 zasoby na darmowym tierze, jeden region."),
            ("03", "Migracje",
             "Start kontenera robi alembic upgrade i wstaje na $PORT — jeden krok."),
        ],
        "callout": "Aplikacja działa na żywo na Render (frankfurt); frontend "
                   "proxuje /api i /ws do backendu — ten sam origin.",
    },
    {
        "kind": "close",
        "eyebrow": "11 / Co dalej",
        "foot": "Podsumowanie",
        "heading": "Co dalej.",
        "panels": [
            ("Teraz",
             ["Trzy grywalne gry z botami RL",
              "Portfel, płatności, retencja",
              "Wdrożone i działające w chmurze"]),
            ("Następne",
             ["Silniejsi agenci pokera (więcej PPO)",
              "Statystyki i tabela wyników graczy",
              "Telemetria i tuning trudności"]),
            ("Później",
             ["Więcej gier i wariantów stołów",
              "Tryb turniejowy",
              "Personalizacja bonusów per gracz"]),
        ],
        "thanks": "Dziękujemy.",
        "thanks_sub": "Demo na żywo: kasyno-jb-frontend.onrender.com",
    },
]

TOTAL = len(SLIDES)


def meta_line(idx: int) -> str:
    """Prawy górny blok meta z numerem strony."""
    return f"Inteligentne Kasyno · Inżynieria oprogramowania<br>Czerwiec 2026 · {idx} / {TOTAL}"


# ════════════════════════════ 1. HTML ══════════════════════════════════════
CSS = """
  :root{
    --ink:#0B0F1E; --slate:#1F293A; --muted:#64748B; --accent:#4F46E5;
    --accent-soft:#EEF0FE; --line:#E6E8F0; --bg:#FFFFFF; --bg-soft:#F7F8FB;
    --stage-w:1280px; --stage-h:720px;
  }
  *{box-sizing:border-box;margin:0;padding:0}
  html,body{height:100%;background:#0B0F1E;}
  body{font-family:'Inter',system-ui,sans-serif;color:var(--slate);
    display:flex;align-items:center;justify-content:center;overflow:hidden;}
  #stage{width:var(--stage-w);height:var(--stage-h);position:relative;
    transform-origin:center center;background:var(--bg);
    box-shadow:0 40px 120px rgba(0,0,0,.55);overflow:hidden;border-radius:6px;}
  .slide{position:absolute;inset:0;padding:54px 72px 56px;display:none;
    flex-direction:column;background:var(--bg);}
  .slide.active{display:flex;}
  .chrome-top{display:flex;justify-content:space-between;align-items:flex-start;}
  .brand{font-weight:800;font-size:13px;letter-spacing:.14em;color:var(--ink);}
  .eyebrow{font-weight:700;font-size:11px;letter-spacing:.16em;color:var(--accent);
    text-transform:uppercase;margin-top:6px;}
  .meta{text-align:right;font-size:11px;color:var(--muted);line-height:1.6;letter-spacing:.04em;}
  .chrome-bottom{position:absolute;left:72px;right:72px;bottom:26px;display:flex;
    justify-content:space-between;font-size:10px;color:var(--muted);
    letter-spacing:.08em;text-transform:uppercase;}
  h1{font-size:56px;line-height:1.04;font-weight:800;color:var(--ink);letter-spacing:-.02em;}
  h2{font-size:42px;line-height:1.08;font-weight:800;color:var(--ink);letter-spacing:-.018em;}
  h2.sm{font-size:34px;}
  .lead{font-size:19px;color:var(--slate);font-weight:400;line-height:1.5;max-width:80%;}
  .sub{font-size:16px;color:var(--muted);line-height:1.5;max-width:80%;}
  .body{flex:1;display:flex;flex-direction:column;justify-content:center;}
  .mono{font-family:'JetBrains Mono',monospace;}
  .cols{display:grid;gap:30px;margin-top:30px;}
  .cols.c2{grid-template-columns:1fr 1fr;}
  .cols.c3{grid-template-columns:repeat(3,1fr);}
  .cols.c4{grid-template-columns:repeat(4,1fr);}
  .blk .bt{font-size:17px;font-weight:700;color:var(--ink);margin-bottom:7px;}
  .blk .bd{font-size:13.5px;color:var(--slate);line-height:1.5;}
  .num{font-family:'JetBrains Mono',monospace;font-weight:700;font-size:13px;
    color:var(--accent);margin-bottom:10px;}
  .panel{background:var(--bg-soft);border:1px solid var(--line);border-radius:14px;padding:24px 26px;}
  .panel .ph{font-size:11px;font-weight:700;letter-spacing:.14em;color:var(--accent);
    text-transform:uppercase;margin-bottom:14px;}
  .li{display:flex;gap:10px;font-size:13.5px;color:var(--slate);line-height:1.5;margin:7px 0;}
  .li .dot{color:var(--accent);font-weight:700;}
  .pipe{display:flex;align-items:stretch;gap:0;margin-top:22px;flex-wrap:wrap;}
  .step{flex:1;min-width:120px;background:var(--bg-soft);border:1px solid var(--line);
    border-radius:12px;padding:16px 14px;text-align:left;}
  .step .si{font-family:'JetBrains Mono',monospace;font-size:11px;color:var(--accent);font-weight:700;}
  .step .sl{font-size:14px;font-weight:700;color:var(--ink);margin-top:6px;line-height:1.2;}
  .step .sd{font-size:11px;color:var(--muted);margin-top:5px;line-height:1.35;}
  .arrow{display:flex;align-items:center;color:var(--accent);font-weight:700;font-size:18px;padding:0 8px;}
  .callout{margin-top:24px;font-size:15px;color:var(--slate);line-height:1.55;
    border-left:3px solid var(--accent);padding-left:16px;max-width:88%;}
  .title-wrap{flex:1;display:flex;flex-direction:column;justify-content:center;}
  .kicker{font-size:12px;font-weight:700;letter-spacing:.18em;color:var(--accent);
    text-transform:uppercase;margin-bottom:22px;}
  .title-foot{font-size:12px;color:var(--muted);letter-spacing:.04em;}
  #pageind{position:fixed;bottom:14px;right:18px;color:#94A3B8;font-size:11px;
    font-family:'JetBrains Mono',monospace;z-index:50;letter-spacing:.1em;}
  #hint{position:fixed;bottom:14px;left:18px;color:#475569;font-size:11px;
    font-family:'JetBrains Mono',monospace;z-index:50;letter-spacing:.06em;}
  @media print{
    body{display:block;background:#fff;}
    #stage{box-shadow:none;border-radius:0;transform:none!important;width:100%;height:100vh;}
    .slide{display:none;}
    .slide.print{display:flex;page-break-after:always;}
    #pageind,#hint{display:none;}
  }
"""


def _chrome_top(eyebrow: str, idx: int, muted_eyebrow: bool = False) -> str:
    cls = ' style="color:var(--muted);font-weight:400"' if muted_eyebrow else ""
    return (f'<div class="chrome-top"><div><div class="brand">{BRAND}</div>'
            f'<div class="eyebrow"{cls}>{eyebrow}</div></div>'
            f'<div class="meta">{meta_line(idx)}</div></div>')


def _chrome_bottom(right: str) -> str:
    return (f'<div class="chrome-bottom"><span>Inteligentne Kasyno · projekt zespołowy</span>'
            f'<span>{right}</span></div>')


def _callout(text: str) -> str:
    return f'<div class="callout">{text}</div>'


def _html_title(s: dict, idx: int) -> str:
    return (
        f'<section class="slide active">'
        f'<div class="chrome-top"><div><div class="brand">{BRAND}</div>'
        f'<div class="eyebrow" style="color:var(--muted);font-weight:400">{s["eyebrow"]}</div></div>'
        f'<div class="meta">{META_TOP}</div></div>'
        f'<div class="title-wrap"><div class="kicker">{s["kicker"]}</div>'
        f'<h1>{s["title"]}</h1>'
        f'<p class="lead" style="margin-top:26px;">{s["lead"]}</p></div>'
        f'<div class="title-foot">{s["foot"]}</div></section>'
    )


def _html_cards(s: dict, idx: int) -> str:
    sub = f'<p class="sub" style="margin-top:14px;">{s["sub"]}</p>' if s.get("sub") else ""
    cards = "".join(
        f'<div class="blk"><div class="num">{n}</div>'
        f'<div class="bt">{t}</div><div class="bd">{b}</div></div>'
        for (n, t, b) in s["items"]
    )
    return (
        f'<section class="slide">{_chrome_top(s["eyebrow"], idx)}'
        f'<div class="body"><h2 class="sm">{s["heading"]}</h2>{sub}'
        f'<div class="cols c{s["cols"]}">{cards}</div>'
        f'{_callout(s["callout"])}</div>{_chrome_bottom(s["foot"])}</section>'
    )


def _html_pipeline(s: dict, idx: int) -> str:
    parts = []
    for i, (si, sl, sd) in enumerate(s["steps"]):
        if i:
            parts.append('<div class="arrow">›</div>')
        parts.append(f'<div class="step"><div class="si">{si}</div>'
                     f'<div class="sl">{sl}</div><div class="sd">{sd}</div></div>')
    sub = f'<p class="sub" style="margin-top:12px;">{s["sub"]}</p>' if s.get("sub") else ""
    return (
        f'<section class="slide">{_chrome_top(s["eyebrow"], idx)}'
        f'<div class="body"><h2 class="sm">{s["heading"]}</h2>{sub}'
        f'<div class="pipe">{"".join(parts)}</div>'
        f'{_callout(s["callout"])}</div>{_chrome_bottom(s["foot"])}</section>'
    )


def _panel_html(ph: str, lis: list[str]) -> str:
    items = "".join(f'<div class="li"><span class="dot">·</span>{x}</div>' for x in lis)
    return f'<div class="panel"><div class="ph">{ph}</div>{items}</div>'


def _html_panels(s: dict, idx: int) -> str:
    sub = f'<p class="sub" style="margin-top:14px;">{s["sub"]}</p>' if s.get("sub") else ""
    panels = "".join(_panel_html(ph, lis) for (ph, _pt, lis) in s["panels"])
    return (
        f'<section class="slide">{_chrome_top(s["eyebrow"], idx)}'
        f'<div class="body"><h2 class="sm">{s["heading"]}</h2>{sub}'
        f'<div class="cols c{s["cols"]}" style="margin-top:24px;">{panels}</div>'
        f'{_callout(s["callout"])}</div>{_chrome_bottom(s["foot"])}</section>'
    )


def _html_close(s: dict, idx: int) -> str:
    panels = "".join(_panel_html(ph, lis) for (ph, lis) in s["panels"])
    return (
        f'<section class="slide">{_chrome_top(s["eyebrow"], idx)}'
        f'<div class="body"><h2 class="sm">{s["heading"]}</h2>'
        f'<div class="cols c3" style="margin-top:26px;">{panels}</div>'
        f'<h2 style="margin-top:40px;">{s["thanks"]}</h2>'
        f'<p class="sub" style="margin-top:10px;">{s["thanks_sub"]}</p>'
        f'</div>{_chrome_bottom(s["foot"])}</section>'
    )


_HTML_RENDER = {
    "title": _html_title, "cards": _html_cards, "pipeline": _html_pipeline,
    "panels": _html_panels, "close": _html_close,
}


def build_html(path: Path) -> None:
    sections = "\n".join(_HTML_RENDER[s["kind"]](s, i + 1) for i, s in enumerate(SLIDES))
    html = f"""<!DOCTYPE html>
<html lang="pl">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Inteligentne Kasyno · Pitch Deck</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&family=JetBrains+Mono:wght@400;500;700&display=swap" rel="stylesheet">
<style>{CSS}</style>
</head>
<body>
<div id="stage">
{sections}
</div>
<div id="pageind">01 / {TOTAL:02d}</div>
<div id="hint">← → · F fullscreen · P print</div>
<script>
  const slides=[...document.querySelectorAll('.slide')];
  let i=0; const ind=document.getElementById('pageind');
  function show(n){{
    i=Math.max(0,Math.min(slides.length-1,n));
    slides.forEach((s,k)=>s.classList.toggle('active',k===i));
    ind.textContent=String(i+1).padStart(2,'0')+' / '+String(slides.length).padStart(2,'0');
  }}
  function fit(){{
    const s=document.getElementById('stage');
    const sc=Math.min(window.innerWidth/1280,window.innerHeight/720);
    s.style.transform='scale('+sc+')';
  }}
  window.addEventListener('resize',fit); fit(); show(0);
  document.addEventListener('keydown',e=>{{
    if(e.key==='ArrowRight'||e.key===' '||e.key==='PageDown'){{show(i+1);e.preventDefault();}}
    else if(e.key==='ArrowLeft'||e.key==='PageUp'){{show(i-1);}}
    else if(e.key==='Home'){{show(0);}}
    else if(e.key==='End'){{show(slides.length-1);}}
    else if(e.key.toLowerCase()==='f'){{if(!document.fullscreenElement)document.documentElement.requestFullscreen();else document.exitFullscreen();}}
    else if(e.key.toLowerCase()==='p'){{
      slides.forEach(s=>s.classList.add('print'));window.print();
      setTimeout(()=>slides.forEach(s=>s.classList.remove('print')),500);
    }}
  }});
  document.getElementById('stage').addEventListener('click',e=>{{
    const x=e.clientX/window.innerWidth; show(x>0.5?i+1:i-1);
  }});
</script>
</body>
</html>
"""
    path.write_text(html, encoding="utf-8")


# ════════════════════════════ 2. PPTX (natywny) ════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W_IN, H_IN = 13.333, 7.5
SX = W_IN / 1280.0
SY = H_IN / 720.0

INTER, MONO = "Inter", "JetBrains Mono"


def IX(px): return Inches(px * SX)
def IY(px): return Inches(px * SY)
def FS(px): return Pt(px * 0.75)


def C(hexstr): return RGBColor.from_string(hexstr)


def _no_autofit(tf):
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(IX(x), IY(y), IX(w), IY(h))
    tf = tb.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    return tf


def add_para(tf, runs, align=PP_ALIGN.LEFT, space_before=0, space_after=0,
             line=1.0, first=False):
    """runs: list of (text, size_px, hexcolor, bold, fontname)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    p.line_spacing = line
    for (text, size, color, bold, font) in runs:
        r = p.add_run()
        r.text = text
        r.font.size = FS(size)
        r.font.bold = bold
        r.font.color.rgb = C(color)
        r.font.name = font
    return p


def rrect(slide, x, y, w, h, fill, line_color=None, radius=0.07):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IX(x), IY(y), IX(w), IY(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C(fill)
    if line_color:
        sh.line.color.rgb = C(line_color)
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def rect(slide, x, y, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IX(x), IY(y), IX(w), IY(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C(fill)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def row(count, gap, x0=72, total_w=1136):
    w = (total_w - gap * (count - 1)) / count
    return [(x0 + i * (w + gap), w) for i in range(count)]


def chrome_top(slide, eyebrow, idx, muted_eyebrow=False):
    tf = textbox(slide, 72, 50, 700, 50)
    add_para(tf, [(BRAND, 13, INK, True, INTER)], first=True)
    eb_color = MUTED if muted_eyebrow else ACCENT
    add_para(tf, [(eyebrow.upper(), 11, eb_color, not muted_eyebrow, INTER)], space_before=3)
    # meta (right)
    tf2 = textbox(slide, 700, 50, 508, 60)
    add_para(tf2, [("Inteligentne Kasyno · Inżynieria oprogramowania", 11, MUTED, False, INTER)],
             align=PP_ALIGN.RIGHT, first=True, line=1.4)
    add_para(tf2, [(f"Czerwiec 2026 · {idx} / {TOTAL}", 11, MUTED, False, INTER)],
             align=PP_ALIGN.RIGHT, line=1.4)


def chrome_bottom(slide, right):
    tf = textbox(slide, 72, 678, 600, 24)
    add_para(tf, [("INTELIGENTNE KASYNO · PROJEKT ZESPOŁOWY", 10, MUTED, False, INTER)], first=True)
    tf2 = textbox(slide, 608, 678, 600, 24)
    add_para(tf2, [(right.upper(), 10, MUTED, False, INTER)], align=PP_ALIGN.RIGHT, first=True)


def callout(slide, text, y):
    rect(slide, 72, y, 3, 60, ACCENT)
    tf = textbox(slide, 90, y, 1010, 70, anchor=MSO_ANCHOR.TOP)
    add_para(tf, [(text, 15, SLATE, False, INTER)], first=True, line=1.45)


def heading(slide, text, y, size=34):
    tf = textbox(slide, 72, y, 1100, size * 1.6)
    add_para(tf, [(text, size, INK, True, INTER)], first=True, line=1.08)


def subtitle(slide, text, y):
    tf = textbox(slide, 72, y, 980, 50)
    add_para(tf, [(text, 16, MUTED, False, INTER)], first=True, line=1.45)


def card_block(slide, x, y, w, h, num, title, body):
    tf = textbox(slide, x, y, w, h)
    add_para(tf, [(num, 13, ACCENT, True, MONO)], first=True, space_after=8)
    add_para(tf, [(title, 17, INK, True, INTER)], space_after=6, line=1.1)
    add_para(tf, [(body, 13.5, SLATE, False, INTER)], line=1.45)


def panel_block(slide, x, y, w, h, ph, lis):
    rrect(slide, x, y, w, h, BG_SOFT, LINE, radius=0.08)
    tf = textbox(slide, x + 24, y + 22, w - 48, h - 40)
    add_para(tf, [(ph.upper(), 11, ACCENT, True, INTER)], first=True, space_after=10)
    for it in lis:
        p = add_para(tf, [("·  ", 13.5, ACCENT, True, INTER),
                          (it, 13.5, SLATE, False, INTER)], space_after=5, line=1.4)


# ── renderery slajdów PPTX ──────────────────────────────────────────────────
def pptx_title(slide, s, idx):
    tf = textbox(slide, 72, 50, 700, 50)
    add_para(tf, [(BRAND, 13, INK, True, INTER)], first=True)
    add_para(tf, [(s["eyebrow"], 11, MUTED, False, INTER)], space_before=3)
    tf2 = textbox(slide, 700, 50, 508, 60)
    add_para(tf2, [("Projekt zespołowy · Inżynieria oprogramowania", 11, MUTED, False, INTER)],
             align=PP_ALIGN.RIGHT, first=True, line=1.4)
    add_para(tf2, [("Czerwiec 2026", 11, MUTED, False, INTER)], align=PP_ALIGN.RIGHT, line=1.4)

    tfk = textbox(slide, 72, 230, 900, 30)
    add_para(tfk, [(s["kicker"].upper(), 12, ACCENT, True, INTER)], first=True)

    title_lines = s["title"].replace("<br>", "\n").split("\n")
    tft = textbox(slide, 72, 270, 1080, 180)
    for i, ln in enumerate(title_lines):
        add_para(tft, [(ln, 54, INK, True, INTER)], first=(i == 0), line=1.04)

    tfl = textbox(slide, 72, 470, 980, 120)
    add_para(tfl, [(s["lead"], 19, SLATE, False, INTER)], first=True, line=1.5)

    tff = textbox(slide, 72, 650, 1100, 30)
    add_para(tff, [(s["foot"], 12, MUTED, False, INTER)], first=True)


def pptx_cards(slide, s, idx):
    chrome_top(slide, s["eyebrow"], idx)
    heading(slide, s["heading"], 150)
    has_sub = bool(s.get("sub"))
    if has_sub:
        subtitle(slide, s["sub"], 218)
    cards_y = 290 if has_sub else 250
    cards_h = 210
    cols = s["cols"]
    gap = 26 if cols >= 3 else 34
    for (x, w), (n, t, b) in zip(row(cols, gap), s["items"]):
        card_block(slide, x, cards_y, w, cards_h, n, t, b)
    callout(slide, s["callout"], 540)
    chrome_bottom(slide, s["foot"])


def pptx_pipeline(slide, s, idx):
    chrome_top(slide, s["eyebrow"], idx)
    heading(slide, s["heading"], 150)
    if s.get("sub"):
        subtitle(slide, s["sub"], 218)
    steps = s["steps"]
    n = len(steps)
    aw = 28  # arrow width px
    total_arrows = aw * (n - 1)
    sw = (1136 - total_arrows) / n
    y = 300
    h = 150
    x = 72
    for i, (si, sl, sd) in enumerate(steps):
        if i:
            tfa = textbox(slide, x, y, aw, h, anchor=MSO_ANCHOR.MIDDLE)
            add_para(tfa, [("›", 20, ACCENT, True, INTER)], align=PP_ALIGN.CENTER, first=True)
            x += aw
        rrect(slide, x, y, sw, h, BG_SOFT, LINE, radius=0.1)
        tf = textbox(slide, x + 14, y + 16, sw - 28, h - 28)
        add_para(tf, [(si, 11, ACCENT, True, MONO)], first=True, space_after=6)
        add_para(tf, [(sl, 14, INK, True, INTER)], space_after=5, line=1.15)
        add_para(tf, [(sd, 11, MUTED, False, INTER)], line=1.35)
        x += sw
    callout(slide, s["callout"], 500)
    chrome_bottom(slide, s["foot"])


def pptx_panels(slide, s, idx):
    chrome_top(slide, s["eyebrow"], idx)
    heading(slide, s["heading"], 150)
    has_sub = bool(s.get("sub"))
    if has_sub:
        subtitle(slide, s["sub"], 218)
    py = 270 if has_sub else 250
    cols = s["cols"]
    ph_h = 255 if cols == 2 else 235
    gap = 34 if cols == 2 else 30
    for (x, w), (ph, _pt, lis) in zip(row(cols, gap), s["panels"]):
        panel_block(slide, x, py, w, ph_h, ph, lis)
    callout(slide, s["callout"], py + ph_h + 25)
    chrome_bottom(slide, s["foot"])


def pptx_close(slide, s, idx):
    chrome_top(slide, s["eyebrow"], idx)
    heading(slide, s["heading"], 120)
    py = 190
    ph_h = 215
    for (x, w), (ph, lis) in zip(row(3, 30), s["panels"]):
        panel_block(slide, x, py, w, ph_h, ph, lis)
    tft = textbox(slide, 72, 440, 1000, 70)
    add_para(tft, [(s["thanks"], 42, INK, True, INTER)], first=True)
    tfs = textbox(slide, 72, 520, 1000, 40)
    add_para(tfs, [(s["thanks_sub"], 16, MUTED, False, INTER)], first=True)
    chrome_bottom(slide, s["foot"])


_PPTX_RENDER = {
    "title": pptx_title, "cards": pptx_cards, "pipeline": pptx_pipeline,
    "panels": pptx_panels, "close": pptx_close,
}


def build_pptx_native(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank = prs.slide_layouts[6]
    for i, s in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        # białe tło
        bg = rect(slide, -2, -2, 1284, 724, BG)
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
        _PPTX_RENDER[s["kind"]](slide, s, i + 1)
    prs.save(str(path))


# ── opcjonalna ścieżka: zrzuty z HTML-a (gdy jest playwright + Chromium) ────
def build_pptx_from_screenshots(html_path: Path, pptx_path: Path) -> bool:
    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return False
    shots_dir = html_path.parent / "_slides"
    shots_dir.mkdir(exist_ok=True)
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 720},
                                    device_scale_factor=2)
            page.goto(html_path.resolve().as_uri())
            page.wait_for_timeout(800)
            files = []
            for n in range(TOTAL):
                page.evaluate(f"show({n})")
                page.wait_for_timeout(150)
                fp = shots_dir / f"slide_{n:02d}.png"
                page.locator("#stage").screenshot(path=str(fp))
                files.append(fp)
            browser.close()
    except Exception as e:
        print(f"  [i] zrzuty nieudane ({e}); używam renderu natywnego.")
        return False

    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank = prs.slide_layouts[6]
    for fp in files:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(fp), 0, 0, width=Inches(W_IN), height=Inches(H_IN))
    prs.save(str(pptx_path))
    return True


def main():
    here = Path(__file__).resolve().parent
    html_path = here / "presentation.html"
    pptx_path = here / "presentation.pptx"

    print("→ Generuję HTML...")
    build_html(html_path)
    print(f"  ✓ {html_path}")

    print("→ Konwertuję do PPTX...")
    if build_pptx_from_screenshots(html_path, pptx_path):
        print(f"  ✓ {pptx_path}  (zrzuty z przeglądarki, 1:1)")
    else:
        build_pptx_native(pptx_path)
        print(f"  ✓ {pptx_path}  (render natywny, edytowalny)")

    print(f"\nGotowe: {TOTAL} slajdów.")
    print(f"  HTML : otwórz {html_path.name} w przeglądarce (strzałki / F / P)")
    print(f"  PPTX : {pptx_path.name}")


if __name__ == "__main__":
    main()
